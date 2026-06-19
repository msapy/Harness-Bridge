import os
from html2image import Html2Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def take_screenshots():
    print("Initializing html2image...")
    hti = Html2Image()
    hti.size = (1280, 800)
    
    explorers = {
        "sv_explorer.html": "sv_explorer.png",
        "phone_sv_explorer.html": "phone_sv_explorer.png",
        "multi_sv_explorer.html": "multi_sv_explorer.png"
    }
    
    screenshot_paths = {}
    for html_name, img_name in explorers.items():
        if os.path.exists(html_name):
            print(f"Capturing screenshot of {html_name}...")
            # html2image takes file path
            hti.screenshot(html_file=html_name, save_as=img_name)
            if os.path.exists(img_name):
                print(f"  Successfully saved {img_name}")
                screenshot_paths[html_name] = img_name
            else:
                print(f"  Failed to save {img_name}")
        else:
            print(f"Warning: {html_name} not found.")
            
    return screenshot_paths

def build_presentation(screenshots):
    print("Building presentation...")
    prs = Presentation()
    
    # Set slide dimensions to 16:9 widescreen (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 1. Title Slide
    blank_slide_layout = prs.slide_layouts[6]
    title_slide = prs.slides.add_slide(blank_slide_layout)
    
    # Dark modern background
    bg = title_slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        0, 0, Inches(13.33), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(8, 12, 20)  # matching explorer theme
    bg.line.fill.background()
    
    # Title Text
    txBox = title_slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Bridge SV Mode Shape Explorers"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)  # sky blue
    p.font.name = "Arial"
    
    p2 = tf.add_paragraph()
    p2.text = "Interactive FDD & Singular Value Spectrums Visualizations"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(167, 139, 250)  # purple
    p2.font.name = "Arial"
    
    p3 = tf.add_paragraph()
    p3.text = "\nCompiled Automated Screenshots Report"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(148, 163, 184)  # slate
    p3.font.name = "Arial"

    # Descriptions for slides
    info = {
        "sv_explorer.html": {
            "title": "Standard SV Mode Shape Explorer",
            "desc": "Visualizes the Singular Value Spectrum from accelerometer data, showing peaks corresponding to natural frequencies. Selecting a frequency updates the spline-interpolated deck displacement shape on the right."
        },
        "phone_sv_explorer.html": {
            "title": "Phone SV Mode Shape Explorer",
            "desc": "Adapts the SV Explorer specifically for Phone sensor data measurements. Displays the frequency spectrum and computed cubic-spline shape matching the phone layout on the bridge."
        },
        "multi_sv_explorer.html": {
            "title": "Multi-CSV SV Mode Shape Explorer",
            "desc": "Enables side-by-side comparison of SV curves and mode shapes across multiple CSV files (e.g. different days or measurement files) in a unified interactive view."
        }
    }
    
    # 2. Add slides for each screenshot
    for html_name, img_name in screenshots.items():
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Dark modern background
        bg = slide.shapes.add_shape(
            1,
            0, 0, Inches(13.33), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(8, 12, 20)
        bg.line.fill.background()
        
        # Slide Title
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = info[html_name]["title"]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(96, 165, 250)
        p.font.name = "Arial"
        
        # Slide Description
        txBox_desc = slide.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(11.83), Inches(0.8))
        tf_desc = txBox_desc.text_frame
        tf_desc.word_wrap = True
        p_desc = tf_desc.paragraphs[0]
        p_desc.text = info[html_name]["desc"]
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = RGBColor(203, 213, 225)
        p_desc.font.name = "Arial"
        
        # Screenshot Image (scaled to fit, centered)
        # 1280x800 is 1.6 aspect ratio
        # Box for image: width = 9 inches, height = 5.0 inches (1.8 ratio)
        # Let's fit 8.0w x 5.0h
        left = Inches(2.66)
        top = Inches(2.1)
        width = Inches(8.0)
        height = Inches(5.0)
        
        slide.shapes.add_picture(img_name, left, top, width=width, height=height)
        
        # Subtle border / frame around image
        frame = slide.shapes.add_shape(
            1, left, top, width, height
        )
        frame.fill.background()
        frame.line.color.rgb = RGBColor(56, 189, 248) # sky border
        frame.line.width = Pt(1.5)
        
    out_pptx = "sv_explorers_presentation.pptx"
    prs.save(out_pptx)
    print(f"Presentation saved to: {out_pptx}")

if __name__ == "__main__":
    screenshots = take_screenshots()
    if screenshots:
        build_presentation(screenshots)
        print("Success!")
    else:
        print("No screenshots captured. Cannot build PPTX.")
